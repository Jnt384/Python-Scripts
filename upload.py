import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

def is_valid_image(filename):
    # Check if the file has a valid image extension
    valid_extensions = ['.png', '.jpg', '.jpeg', '.gif', '.bmp']
    return any(filename.lower().endswith(ext) for ext in valid_extensions)

def add_screenshot_to_ppt(presentation, screenshot_path):
    try:
        # Add a new slide to the PowerPoint presentation
        slide_layout = presentation.slide_layouts[5]  # Use the layout that suits your needs
        slide = presentation.slides.add_slide(slide_layout)

        # Calculate the position to add the image on the slide
        left = Inches(1)  # Adjust this value as needed
        top = Inches(1)   # Adjust this value as needed

        # Add the image to the PowerPoint presentation
        picture = slide.shapes.add_picture(screenshot_path, left, top, width=Inches(6), height=Inches(4))

    except Exception as e:
        print(f"An error occurred while adding screenshot: {e}")

def main():   
    try:
        ppt_folder_path = r'C:\Users\jt84h\OneDrive\Documents\reactjs\Python\screenshots\Screenshot_to_PPT'
        ppt_file_name = "CEH.pptx"

        if os.path.exists(os.path.join(ppt_folder_path, ppt_file_name)):
            presentation = Presentation(os.path.join(ppt_folder_path, ppt_file_name))
        else:
            presentation = Presentation()

        for file_name in os.listdir(ppt_folder_path):
            file_path = os.path.join(ppt_folder_path, file_name)

            # Check if the file is a valid image
            if is_valid_image(file_name):
                add_screenshot_to_ppt(presentation, file_path)
                print(f"Screenshot added: {file_name}")

        presentation.save(os.path.join(ppt_folder_path, ppt_file_name))
        print(f"All new screenshots added to {ppt_file_name}")

    except Exception as e:
        print(f"An error occurred: {e}")

if __name__ == "__main__":
    main()

